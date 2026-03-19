"""
One-off script: build a single PowerPoint slide with title "W3LPL Station Architecture"
and the RF system architecture diagram image. Run from repo root or from assets/.
"""
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    raise SystemExit("Install python-pptx: pip install python-pptx")

# Paths: assume script is in repo/assets/
SCRIPT_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRIPT_DIR.parent
IMAGE_PATH = SCRIPT_DIR / "rf-system-architecture-cleaned.png"
OUTPUT_PATH = SCRIPT_DIR / "W3LPL_Station_Architecture.pptx"

# Fallback: image might be in .cursor project assets
if not IMAGE_PATH.exists():
    CURSOR_ASSETS = REPO_ROOT.parent.parent / ".cursor" / "projects" / "c-Users-mbdev-OneDrive-Desktop-Repos-Contest-Log-Analyzer" / "assets"
    alt = CURSOR_ASSETS / "rf-system-architecture-cleaned.png"
    if alt.exists():
        IMAGE_PATH = alt

if not IMAGE_PATH.exists():
    raise FileNotFoundError(
        f"Diagram image not found at {SCRIPT_DIR / 'rf-system-architecture-cleaned.png'} "
        f"or at {IMAGE_PATH}. Ensure assets/rf-system-architecture-cleaned.png exists."
    )

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Blank layout
blank = prs.slide_layouts[6]  # Blank
slide = prs.slides.add_slide(blank)

# Title at top
from pptx.dml.color import RgbColor
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.8))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "W3LPL Station Architecture"
p.font.size = Pt(32)
p.font.bold = True
p.font.name = "Calibri"
try:
    p.font.color.rgb = RgbColor(0x1F, 0x4E, 0x79)
except Exception:
    pass

# Diagram image: max width 12", preserve aspect ratio, center below title
slide.shapes.add_picture(
    str(IMAGE_PATH),
    Inches(0.5),
    Inches(1.25),
    width=Inches(12.333),
)

prs.save(str(OUTPUT_PATH))
print(f"Saved: {OUTPUT_PATH}")
